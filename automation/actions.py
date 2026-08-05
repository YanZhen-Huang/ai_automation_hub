"""动作库：会议准备流程中可执行的动作。执行结果写入 action_logs。"""

import datetime
import subprocess
import webbrowser
from pathlib import Path

from core.config import DATA_DIR
from core.events import bus
from core.logger import get_logger
from storage import db

log = get_logger("actions")
OUT_DIR = DATA_DIR / "out"
OUT_DIR.mkdir(parents=True, exist_ok=True)


def _log(item_id, action, status, message=""):
    db.add_action_log(item_id, action, status, message)
    bus().publish("action.logged", {"item_id": item_id, "action": action,
                                    "status": status, "message": message})


def sanitize_filename(name, max_len=80):
    """清洗 Windows 非法文件名字符。"""
    import re as _re
    s = _re.sub(r'[\\/:*?"<>|]', "_", str(name or "untitled"))
    s = s.strip().strip(".")
    if len(s) > max_len:
        s = s[:max_len]
    return s or "untitled"


def md_to_pdf(md_text, out_path):
    """markdown → PDF（fpdf2 + 中文字体）。返回路径或 None。"""
    try:
        from fpdf import FPDF
    except ImportError:
        return None
    import re as _re
    FONT = r"C:\Windows\Fonts\simhei.ttf"
    pdf = FPDF()
    pdf.add_font("hei", "", FONT)
    pdf.add_page()
    pdf.set_margins(14, 12, 14)
    cw = pdf.w - 28
    for raw in str(md_text or "").splitlines():
        s = raw.strip()
        if not s:
            continue
        if s.startswith("# "):
            pdf.set_font("hei", size=16)
            pdf.multi_cell(cw, 9, s[2:])
        elif s.startswith("## "):
            pdf.set_font("hei", size=13)
            pdf.multi_cell(cw, 8, s[3:])
        elif s.startswith("- "):
            pdf.set_font("hei", size=10)
            pdf.multi_cell(cw, 6, "· " + s[2:])
        elif s.startswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if not all(_re.fullmatch(r":?-{2,}:?", c) for c in cells):
                pdf.set_font("hei", size=9)
                pdf.cell(cw, 6, " | ".join(cells))
                pdf.ln()
        else:
            pdf.set_font("hei", size=10)
            pdf.multi_cell(cw, 6, s)
    pdf.output(str(out_path))
    return str(out_path)


def notify(title, message):
    """系统通知（由桌面/Web 端展示）。"""
    bus().publish("notification", {"title": title, "message": message})
    log.info("通知 [%s] %s", title, message)
    return True


def export_text(title, content, ext="md"):
    """导出文本/材料文件，返回路径。"""
    import datetime
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe = sanitize_filename(title)
    path = OUT_DIR / f"{stamp}_{safe}.{ext}"
    path.write_text(content or "", encoding="utf-8")
    log.info("已导出: %s", path)
    return str(path)


def generate_ppt(title, points, path=None):
    """用要点生成 PPT 初稿，返回路径。"""
    try:
        from pptx import Presentation
    except ImportError:
        return None
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title or "汇报材料"
    for p in (points or [])[:6]:
        body = slide.placeholders[1]
        tf = body.text_frame
        para = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
        para.text = p
    if path is None:
        import datetime
        stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        path = str(OUT_DIR / f"{stamp}_{sanitize_filename(title)}.pptx")
    prs.save(path)
    log.info("已生成 PPT: %s", path)
    return path


def open_url(url):
    webbrowser.open(url)
    return True


def run_command(cmd):
    try:
        if isinstance(cmd, str):
            cmd = cmd.split()
        subprocess.Popen(cmd, shell=False)
        return True
    except Exception:
        log.exception("命令执行失败: %s", cmd)
        return False


def _get_wechat_uia():
    try:
        import uiautomation as auto
        return auto
    except Exception:
        return None


def _find_wechat_window(auto):
    root = auto.GetRootControl()
    for w in root.GetChildren():
        try:
            if not w.IsWindow:
                continue
            name = w.Name or ""
            cls = w.ClassName or ""
            if ("微信" in name) or ("WeChat" in cls):
                return w
        except Exception:
            continue
    return None


def _collect_edits(win, limit=6):
    edits = []

    def walk(c, depth):
        if len(edits) >= limit or depth > 8:
            return
        try:
            children = c.GetChildren()
        except Exception:
            return
        for ch in children:
            try:
                if ch.ControlTypeName == "EditControl":
                    edits.append(ch)
            except Exception:
                pass
            walk(ch, depth + 1)

    walk(win, 0)
    return edits


def _find_conv_item(win, target):
    best = None

    def walk(c, depth):
        nonlocal best
        if best is not None or depth > 8:
            return
        try:
            children = c.GetChildren()
        except Exception:
            return
        for ch in children:
            try:
                name = (ch.Name or "").strip()
                if name and target and target in name:
                    best = ch
                    return
            except Exception:
                pass
            walk(ch, depth + 1)

    walk(win, 0)
    return best


def send_wechat_text(text, target=""):
    """通过 UIA 自动向微信目标(群/联系人)发送文本。失败返回 False。"""
    if not (text and target):
        return False
    from core.config import CONFIG
    if not CONFIG.get("wechat", {}).get("send_enabled", False):
        return False
    import time
    auto = _get_wechat_uia()
    if auto is None:
        return False
    win = _find_wechat_window(auto)
    if win is None:
        return False
    try:
        _deadline = time.time() + 30

        def _timeout():
            return time.time() > _deadline

        edits = _collect_edits(win)
        if not edits:
            return False
        search = edits[0]
        search.Click()
        time.sleep(0.3)
        if _timeout():
            return False
        search.SendKeys("{Ctrl}a")
        search.SendKeys(target)
        time.sleep(int(CONFIG.get("wechat", {}).get("send_delay", 3)))
        if _timeout():
            return False
        conv = _find_conv_item(win, target)
        if conv is None:
            log.warning("未找到会话: %s", target)
            return False
        conv.Click()
        time.sleep(0.6)
        if _timeout():
            return False
        edits2 = _collect_edits(win)
        if not edits2:
            return False
        inbox = edits2[-1]
        inbox.Click()
        time.sleep(0.3)
        if _timeout():
            return False
        inbox.SendKeys(text)
        time.sleep(0.3)
        inbox.SendKeys("{Enter}")
        return True
    except Exception:
        log.exception("微信自动发送失败")
        return False


def call_phone(item_id, action, text, number=""):
    """提交打电话任务给手机端联动。返回 task_id 或 None。"""
    from automation.phone_link import PhoneLink
    return PhoneLink.instance().submit_call(item_id, action, text, number)


def send_wechat_file(path, target=""):
    """通过微信 UIA 自动向目标发送文件（用于交付印刷材料）。"""
    if not (path and target):
        return False
    from core.config import CONFIG
    if not CONFIG.get("wechat", {}).get("send_enabled", False):
        return False
    if not Path(path).exists():
        return False
    import time
    auto = _get_wechat_uia()
    if auto is None:
        return False
    win = _find_wechat_window(auto)
    if win is None:
        return False
    try:
        edits = _collect_edits(win)
        if not edits:
            return False
        search = edits[0]
        search.Click()
        time.sleep(0.3)
        search.SendKeys("{Ctrl}a")
        search.SendKeys(target)
        time.sleep(int(CONFIG.get("wechat", {}).get("send_delay", 3)))
        conv = _find_conv_item(win, target)
        if conv is None:
            log.warning("未找到文件发送目标: %s", target)
            return False
        conv.Click()
        time.sleep(0.6)
        # 点击"+"附件入口
        plus = _find_button_by_name(win, ["文件", "附件", "更多", "+"])
        if plus is None:
            return False
        plus.Click()
        time.sleep(0.8)
        if not _open_file_dialog(path):
            return False
        time.sleep(0.6)
        # 回车发送
        _find_send_enter(win)
        return True
    except Exception:
        log.exception("微信发送文件失败")
        return False


def _find_button_by_name(win, names):
    best = None

    def walk(c, depth):
        nonlocal best
        if best is not None or depth > 6:
            return
        try:
            children = c.GetChildren()
        except Exception:
            return
        for ch in children:
            try:
                n = (ch.Name or "").strip()
                if any(k in n for k in names):
                    best = ch
                    return
            except Exception:
                pass
            walk(ch, depth + 1)

    walk(win, 0)
    return best


def _open_file_dialog(path):
    """在 Windows 打开对话框(#32770)中输入路径并确认。"""
    auto = _get_wechat_uia()
    if auto is None:
        return False
    root = auto.GetRootControl()
    for w in root.GetChildren():
        try:
            if w.ClassName == "#32770" and ("打开" in (w.Name or "")
                                            or "选择" in (w.Name or "")):
                edit = w.Control(searchDepth=3, ControlType="Edit")
                if edit is not None:
                    edit.Click()
                    edit.SendKeys("{Ctrl}a")
                    edit.SendKeys(path)
                    ok = w.Control(searchDepth=3, Name="打开")
                    if ok is not None:
                        ok.Click()
                        return True
        except Exception:
            continue
    return False


def _find_send_enter(win):
    """找到发送按钮并点击；找不到则尝试回车。"""
    try:
        btn = _find_button_by_name(win, ["发送"])
        if btn is not None:
            btn.Click()
            return True
    except Exception:
        pass
    try:
        win.SendKeys("{Enter}")
        return True
    except Exception:
        return False


def execute(action, item_id, **kwargs):
    """统一执行入口，带日志记录。"""
    result = None
    try:
        if action == "notify":
            result = notify(kwargs.get("title", ""), kwargs.get("message", ""))
        elif action == "export":
            result = export_text(kwargs.get("title", ""), kwargs.get("content", ""))
        elif action == "ppt":
            result = generate_ppt(kwargs.get("title", ""), kwargs.get("points", []))
        elif action == "open_url":
            result = open_url(kwargs.get("url", ""))
        elif action == "run_command":
            result = run_command(kwargs.get("cmd", ""))
        elif action == "wechat":
            result = send_wechat_text(kwargs.get("text", ""),
                                      kwargs.get("target", ""))
        elif action == "call":
            result = call_phone(item_id, kwargs.get("name", ""),
                                kwargs.get("text", ""), kwargs.get("number", ""))
        else:
            raise ValueError(f"未知动作: {action}")
        _log(item_id, action, "success", str(result))
        return result
    except Exception as e:
        _log(item_id, action, "error", str(e))
        log.exception("动作执行失败: %s", action)
        return None
