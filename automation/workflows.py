"""固定工作流：会议准备。

一阶段（并行，除特殊要求外）：订会议室/茶水/服务/音响设施、备材料转PPT+印刷、人工定时间、微信通知、定签字表。
二阶段（串行）：定与会人员 → 打电话排桌牌 → 人工审查 → 提醒接待。
"""

import json
import re
from pathlib import Path

from ai.extract import extract, generate_materials
from automation import actions, phrases
from core.config import CONFIG
from core.events import bus
from core.logger import get_logger
from storage import db

OUT_DIR = actions.OUT_DIR

log = get_logger("workflow")

PHASE1 = [
    {"code": "call_room", "name": "打电话要会议室"},
    {"code": "call_tea", "name": "打电话要茶水"},
    {"code": "call_service", "name": "打电话要服务"},
    {"code": "call_facilities", "name": "打电话要音响/投影等配套设施"},
    {"code": "prep_materials", "name": "准备汇报材料并转PPT后安排印刷"},
    {"code": "confirm_time", "name": "确定会议时间（人工）"},
    {"code": "notify_wechat", "name": "微信下通知"},
    {"code": "confirm_sign_list", "name": "确定签字表"},
]
PHASE2 = [
    {"code": "confirm_attendees", "name": "确定与会人员"},
    {"code": "call_table_card", "name": "打电话安排桌牌"},
    {"code": "manual_review", "name": "人工审查确认"},
    {"code": "remind_reception", "name": "提醒人工接待"},
]

CALL_CODES = {"call_room", "call_tea", "call_service", "call_facilities",
              "call_table_card"}
HUMAN_CODES = {"confirm_time", "confirm_sign_list", "confirm_attendees",
               "manual_review"}
# 需审批后真正执行的自动化动作
AUTO_CODES = CALL_CODES | {"prep_materials", "notify_wechat"}


def setup_meeting(title, start_time=None):
    """创建会议并挂载固定工作流动作项。返回 meeting_id。"""
    mid = db.create_meeting(title, start_time)
    items = []
    for i, it in enumerate(PHASE1):
        items.append({"phase": 1, **it, "order_index": i + 1})
    for i, it in enumerate(PHASE2):
        items.append({"phase": 2, **it, "order_index": i + 1})
    db.add_prep_items(mid, items)
    return mid


def run_phase1(meeting_id):
    """并行触发一阶段：所有待执行动作先进入"待审批"，由用户批准后执行。
    同时激活二阶段首项（待审批）。"""
    for item in db.list_prep_items(meeting_id):
        if item["phase"] == 1 and item["status"] == "pending":
            request_approval(item)
    _advance_phase2(meeting_id)


def request_approval(item):
    """将动作项置为"待审批"并通知用户。"""
    db.update_prep_item(item["id"], status="waiting", result="待审批")
    bus().publish("approval.requested",
                  {"item_id": item["id"], "name": item["name"],
                   "meeting_id": item["meeting_id"]})


def approve(item_id):
    """用户审批通过后执行；人工确认类动作则直接标记完成。
    会议室/桌牌动作若缺少必要信息（会议室/人员名单），先请求用户补充。"""
    item = db.get_prep_item(item_id)
    if item is None or item["status"] != "waiting":
        return
    meeting = db.get_meeting(item["meeting_id"]) or {}
    if item["code"] == "call_room" and not (meeting.get("room") or "").strip():
        bus().publish("room.select.requested",
                      {"item_id": item_id, "meeting_id": item["meeting_id"],
                       "name": item["name"]})
        return
    if item["code"] in ("call_table_card", "prep_materials") \
            and not (meeting.get("attendees") or "").strip():
        # 桌牌 / 印刷份数都依赖与会人员：人员未定先请求确定，确定后自动继续
        bus().publish("attendees.requested",
                      {"item_id": item_id, "meeting_id": item["meeting_id"],
                       "name": item["name"]})
        return
    if item["code"] in HUMAN_CODES:
        _mark_done(item_id, "已确认")
        return
    if item["code"] == "prep_materials":
        # 双线：auto_materials=1 自动全流程；=0 生成文件后人工发送
        if meeting.get("auto_materials", 1):
            execute_item(item)
        else:
            gen_files(item["meeting_id"])
            _mark_done(item["id"], "材料已生成，请到会议详情手动选择文件发送")
            actions.notify("材料已生成（人工发送）",
                           "请到「材料/文件」区选择文件发送给印刷方")
        return
    execute_item(item)


def set_room(meeting_id, room):
    """用户选定会议室后，继续执行预订动作。"""
    db.update_meeting(meeting_id, room=room)
    _resume_waiting(meeting_id, "call_room")


def set_attendees(meeting_id, attendees, source="manual"):
    """用户确定与会人员后，继续执行依赖人员的动作（桌牌 / 印刷份数）。"""
    db.update_meeting(meeting_id, attendees=attendees, attendee_source=source)
    _resume_waiting(meeting_id, "call_table_card")
    _resume_waiting(meeting_id, "prep_materials")


def _resume_waiting(meeting_id, code):
    for it in db.list_prep_items(meeting_id):
        if it["code"] == code and it["status"] == "waiting":
            approve(it["id"])


def candidates_for_room(meeting_id=None):
    """候选会议室 = 会议室库 - AI 提炼的不可用项（去重）。"""
    from ai.extract import extract_rooms
    infos = db.list_info_items(meeting_id) or db.list_info_items()
    unavailable = set(extract_rooms([i["content"] for i in infos]))
    seen, out = set(), []
    for r in db.list_rooms():
        if r["status"] != "active" or r["name"] in unavailable:
            continue
        if r["name"] in seen:
            continue
        seen.add(r["name"])
        out.append(r["name"])
    return out


def execute_item(item):
    code = item["code"]
    meeting = db.get_meeting(item["meeting_id"]) or {}
    try:
        if code in CALL_CODES:
            db.update_prep_item(item["id"], status="running")
            text = phrases.fill(code, meeting)
            actions.call_phone(item["id"], item["name"], text, db.get_number(code))

        elif code == "prep_materials":
            db.update_prep_item(item["id"], status="running")
            _prep_materials(item, meeting)

        elif code == "notify_wechat":
            text = _notify_text(meeting)
            target = _notify_target(meeting)
            ok = actions.send_wechat_text(text, target=target)
            if ok:
                result = f"已自动发送到「{target}」"
            else:
                actions.notify("微信通知（请手动发送）", f"目标:{target}\n{text}")
                result = f"未自动发送（目标:{target}），已生成通知文本"
            _mark_done(item["id"], result)

        elif code in HUMAN_CODES:
            db.update_prep_item(item["id"], status="waiting")
            actions.notify("会议准备 · 需人工确认", item["name"])

        elif code == "remind_reception":
            actions.notify("会议准备", "请提醒人工接待与会人员。")
            _mark_done(item["id"], "已提醒接待")

        else:
            _mark_done(item["id"], "")
    except Exception:
        log.exception("动作执行异常: %s", item["name"])
        db.update_prep_item(item["id"], status="waiting", result="执行出错，请重试")


def mark_done(item_id, result=""):
    """人工确认/完成某动作项。"""
    item = db.get_prep_item(item_id)
    if item is None:
        return
    _mark_done(item_id, result)


def _mark_done(item_id, result):
    db.update_prep_item(item_id, status="done", result=result)
    item = db.get_prep_item(item_id)
    if item and item["phase"] == 2:
        _advance_phase2(item["meeting_id"])
    bus().publish("prep.updated", {"item_id": item_id})


def _advance_phase2(meeting_id):
    """串行推进：找到第一个未完成项，若前一项已完成则进入"待审批"。"""
    items = [i for i in db.list_prep_items(meeting_id) if i["phase"] == 2]
    for idx, item in enumerate(items):
        if item["status"] == "done":
            continue
        if idx == 0 or items[idx - 1]["status"] == "done":
            if item["status"] == "pending":
                request_approval(item)
        break


def on_phone_result(payload):
    item_id = payload.get("item_id")
    if not item_id:
        return
    if payload.get("ok"):
        _mark_done(item_id, payload.get("confirmed") or "已确认")
    else:
        # 拨打失败：回到待审批，用户可重试
        db.update_prep_item(item_id, status="waiting",
                            result=payload.get("confirmed") or "拨打失败，可重试")
        bus().publish("prep.updated", {"item_id": item_id})


def _notify_text(meeting):
    parts = [f"会议通知：{meeting.get('title') or ''}"]
    if meeting.get("start_time"):
        parts.append(f"时间：{meeting['start_time']}")
    if meeting.get("room"):
        parts.append(f"地点：{meeting['room']}")
    parts.append("请准时参加。")
    return "\n".join(p for p in parts if p)


def _notify_target(meeting):
    return (meeting.get("notify_target") or "").strip() or \
        CONFIG.get("wechat", {}).get("notify_target", "") or "未指定"


def _attendee_count(meeting):
    att = (meeting.get("attendees") or "").strip()
    if not att:
        return 0
    return len([a for a in re.split(r"[，,、\s]+", att) if a.strip()])


def gen_materials(meeting_id):
    """AI 生成结构化材料并存储到会议。返回材料 dict。"""
    meeting = db.get_meeting(meeting_id) or {}
    infos = db.list_info_items(meeting_id) or []
    texts = [i["content"] for i in infos[:50]]
    m = generate_materials(meeting.get("title") or "汇报材料", texts, meeting)
    db.update_meeting(meeting_id, materials=json.dumps(m, ensure_ascii=False))
    return m


def get_materials(meeting_id):
    meeting = db.get_meeting(meeting_id) or {}
    if meeting.get("materials"):
        try:
            return json.loads(meeting["materials"])
        except Exception:
            return None
    return None


def gen_files(meeting_id):
    """基于材料生成 PPT/md/PDF，按会议分组登记。返回 [(ftype, path)]。"""
    meeting = db.get_meeting(meeting_id) or {}
    m = get_materials(meeting_id) or gen_materials(meeting_id)
    title = meeting.get("title") or "汇报材料"
    out_dir = OUT_DIR / actions.sanitize_filename(f"{meeting_id}_{title}")
    ppt = _make_ppt(title, m, out_dir)
    doc_txt = _make_doc(title, m)
    doc = out_dir / f"{actions.sanitize_filename(title)}.md"
    doc.write_text(doc_txt, encoding="utf-8")
    pdf = actions.md_to_pdf(doc_txt, out_dir / f"{actions.sanitize_filename(title)}.pdf")
    files = []
    if ppt:
        files.append(("pptx", ppt))
    files.append(("md", str(doc)))
    if pdf:
        files.append(("pdf", pdf))
    for ftype, p in files:
        db.add_meeting_file(meeting_id, Path(p).name, p, ftype)
    return files


def send_files_to_print(meeting_id):
    """把已生成文件发给印刷微信联系人，返回是否成功。"""
    files = db.list_meeting_files(meeting_id)
    print_target = CONFIG.get("wechat", {}).get("print_target", "") or ""
    if not print_target:
        return False
    for pref in ("pdf", "pptx", "md"):
        for f in files:
            if f["ftype"] == pref and not f["sent"]:
                if actions.send_wechat_file(f["path"], print_target):
                    db.update_file_sent(f["id"], 1)
                    return True
    return False


def _prep_materials(item, meeting):
    """自动模式：生成材料+文件 → 自动发印刷 → 电话安排印刷（份数=人数+5）。"""
    files = gen_files(item["meeting_id"])
    sent = send_files_to_print(item["meeting_id"])
    actions.notify("汇报材料已生成",
                   f"文件数：{len(files)}\n已自动发送印刷：{'是' if sent else '否'}")

    count = _attendee_count(meeting) + 5
    first = files[0][1] if files else ""
    text = phrases.fill("call_print", meeting, {"file": first, "count": str(count)})
    actions.call_phone(item["id"], "安排印刷", text, db.get_number("call_print"))


def _make_ppt(title, materials, out_dir):
    """用结构化材料生成 16:9 深蓝科技风 PPT：封面/目录/章节(含表格)/总结/页码。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
    except ImportError:
        return None
    import datetime

    INK = "0A0E1A"; PANEL = "111827"; CYAN = "38BDF8"
    CYAN2 = "4DD0E1"; TXT = "C9D2DD"; GOLD = "D4AF7A"; RED = "E25C6A"

    def _set_run(run, size=18, bold=False, color=TXT):
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", "Microsoft YaHei")

    def _txt(slide, x, y, w, h, text, size=18, bold=False, color=TXT,
             align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = ln
            _set_run(r, size=size, bold=bold, color=color)
        return tb

    def _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(color)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    cover = materials.get("cover") or {}
    chapters = materials.get("chapters") or []
    agenda = materials.get("agenda") or []
    n = len(chapters) + 4 if (agenda or materials.get("summary")) else len(chapters) + 3

    def _footer(slide, page):
        _txt(slide, 11.7, 7.05, 1.4, 0.4, str(page), size=10, color="5A5A66",
             align=PP_ALIGN.RIGHT)
        _txt(slide, 0.4, 7.05, 5, 0.4, title, size=9, color="5A5A66")

    # 封面
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, INK)
    _rect(s, 0, 0, 13.333, 0.12, CYAN2)
    _rect(s, 11.0, 0.6, 1.2, 1.2, RED, MSO_SHAPE.OVAL)
    _rect(s, 10.4, 1.0, 2.2, 1.0, CYAN, MSO_SHAPE.ISOSCELES_TRIANGLE)
    _rect(s, 10.6, 6.2, 2.3, 0.06, GOLD)
    _txt(s, 0.9, 1.4, 10, 1.0, cover.get("title") or title, size=40, bold=True,
         color=CYAN2)
    _txt(s, 0.9, 2.5, 10, 0.6, cover.get("subtitle") or "汇报材料", size=20,
         color=TXT)
    meta = " | ".join(x for x in [
        cover.get("meeting_time"), cover.get("location"),
        cover.get("presenter"), cover.get("department")] if x)
    if meta:
        _txt(s, 0.9, 3.3, 10, 0.5, meta, size=14, color=GOLD)
    _footer(s, 1)

    # 目录
    if agenda:
        s = prs.slides.add_slide(blank)
        _rect(s, 0, 0, 13.333, 7.5, INK)
        _txt(s, 0.9, 0.5, 6, 0.7, "目录", size=28, bold=True, color=CYAN)
        y = 1.5
        for i, a in enumerate(agenda, 1):
            _txt(s, 1.2, y, 10, 0.5, f"{i}. {a}", size=16, color=TXT)
            y += 0.55
        _footer(s, 2)

    # 章节
    for idx, ch in enumerate(chapters, 1):
        s = prs.slides.add_slide(blank)
        _rect(s, 0, 0, 13.333, 7.5, INK)
        _rect(s, 0, 0, 0.14, 7.5, CYAN)
        _txt(s, 0.8, 0.5, 11, 0.7, ch.get("heading") or "章节",
             size=26, bold=True, color=CYAN2)
        y = 1.4
        for p in (ch.get("points") or [])[:7]:
            _txt(s, 1.1, y, 11.5, 0.5, f"▪  {p}", size=14, color=TXT)
            y += 0.5
        tb = ch.get("table")
        if tb and tb.get("cols"):
            try:
                cols = tb["cols"]
                rows = tb.get("rows") or []
                tbl_shape = s.shapes.add_table(len(rows) + 1, len(cols),
                                               Inches(1.2), Inches(y),
                                               Inches(10.5), Inches(0.4 * (len(rows) + 1)))
                tbl = tbl_shape.table
                for j, c in enumerate(cols):
                    cell = tbl.cell(0, j)
                    cell.text = str(c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(PANEL)
                for i, row in enumerate(rows, 1):
                    for j, v in enumerate(row):
                        if j < len(cols):
                            cell = tbl.cell(i, j)
                            cell.text = str(v)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(INK)
            except Exception:
                log.exception("PPT 表格渲染失败")
        _footer(s, idx + 2)

    # 总结
    conclusion = materials.get("conclusion")
    if conclusion:
        s = prs.slides.add_slide(blank)
        _rect(s, 0, 0, 13.333, 7.5, INK)
        _rect(s, 0, 0, 0.14, 7.5, GOLD)
        _txt(s, 0.8, 0.5, 6, 0.7, "总结", size=28, bold=True, color=GOLD)
        _txt(s, 1.2, 1.5, 11, 4, conclusion, size=16, color=TXT)
        _footer(s, n)

    out_dir.mkdir(parents=True, exist_ok=True)
    base = actions.sanitize_filename(title)
    path = out_dir / f"{base}.pptx"
    i = 1
    while path.exists():
        path = out_dir / f"{base}_{i}.pptx"
        i += 1
    prs.save(str(path))
    return str(path)


def _make_doc(title, materials):
    """生成汇报材料 markdown 文档。"""
    import datetime
    cover = materials.get("cover") or {}
    lines = [f"# {cover.get('title') or title}", ""]
    meta = " | ".join(x for x in [
        cover.get("meeting_time"), cover.get("location"),
        cover.get("presenter"), cover.get("department")] if x)
    if meta:
        lines += [f"**{meta}**", ""]
    if materials.get("summary"):
        lines += [f"**摘要**：{materials['summary']}", ""]
    if materials.get("agenda"):
        lines += ["## 议程", ""]
        lines += [f"{i}. {a}" for i, a in enumerate(materials["agenda"], 1)]
        lines.append("")
    for ch in (materials.get("chapters") or []):
        lines += [f"## {ch.get('heading', '')}", ""]
        lines += [f"- {p}" for p in (ch.get("points") or [])]
        tb = ch.get("table")
        if tb and tb.get("cols"):
            lines.append("")
            lines.append("| " + " | ".join(str(c) for c in tb["cols"]) + " |")
            lines.append("|" + "---|" * len(tb["cols"]))
            for row in (tb.get("rows") or []):
                lines.append("| " + " | ".join(str(v) for v in row[:len(tb["cols"])]) + " |")
        lines.append("")
    if materials.get("conclusion"):
        lines += ["## 总结", "", materials["conclusion"], ""]
    return "\n".join(lines)


def on_phone_timeout(payload):
    item_id = payload.get("item_id")
    if item_id:
        db.update_prep_item(item_id, status="waiting", result="手机未回传，已超时，可重试")
        bus().publish("prep.updated", {"item_id": item_id})


bus().subscribe("phone.result", on_phone_result)
bus().subscribe("phone.timeout", on_phone_timeout)
